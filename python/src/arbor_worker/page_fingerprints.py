from __future__ import annotations

import hashlib
import shutil
import subprocess
import tempfile
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation

from arbor_worker.course_manifest import FingerprintKind
from arbor_worker.prepare import PrepareError
from arbor_worker.prepare.pdf import render_pdf_to_images
from arbor_worker.prepare.pptx import convert_pptx_to_pdf, extract_pptx_text
from arbor_worker.settings import WorkerSettings, default_settings
from arbor_worker.sources import classify


@dataclass(frozen=True)
class PageFingerprintResult:
    kind: FingerprintKind
    fingerprints: list[str]


def fingerprint_source(
    source: Path,
    settings: WorkerSettings | None = None,
    *,
    runner=None,
    which=None,
) -> PageFingerprintResult:
    source = Path(source)
    settings = settings or default_settings()
    source_type = classify(source)
    if source_type == "pdf":
        return _fingerprint_pdf(source, settings)
    if source_type == "pptx":
        return _fingerprint_pptx(source, settings, runner=runner, which=which)
    raise PrepareError(f"Unsupported source type: {source.suffix}")


def _digest(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _normalize_text(text: str) -> str:
    return " ".join(text.split())


def _nonspace_len(text: str) -> int:
    return len("".join(text.split()))


def _hash_images(images: list[Path]) -> list[str]:
    return [_digest(path.read_bytes()) for path in images]


def _fingerprint_pdf(source: Path, settings: WorkerSettings) -> PageFingerprintResult:
    with tempfile.TemporaryDirectory(prefix="arbor-fp-") as tmp:
        images = render_pdf_to_images(source, Path(tmp), dpi=settings.pdf_render_dpi)
        return PageFingerprintResult("pdf_image", _hash_images(images))


def _slide_texts(source: Path) -> list[str]:
    slides: list[str] = []
    for slide in Presentation(str(source)).slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        chunks.append(line)
        slides.append("\n".join(chunks))
    return slides


def _fingerprint_pptx(
    source: Path,
    settings: WorkerSettings,
    *,
    runner=None,
    which=None,
) -> PageFingerprintResult:
    combined = extract_pptx_text(source)
    if _nonspace_len(combined) >= settings.pptx_min_chars:
        fingerprints = [
            _digest(_normalize_text(text).encode()) for text in _slide_texts(source)
        ]
        return PageFingerprintResult("pptx_text", fingerprints)

    with tempfile.TemporaryDirectory(prefix="arbor-fp-") as tmp:
        out_dir = Path(tmp)
        pdf = convert_pptx_to_pdf(
            source,
            out_dir,
            runner=runner or subprocess.run,
            which=which or shutil.which,
        )
        images = render_pdf_to_images(pdf, out_dir, dpi=settings.pdf_render_dpi)
        return PageFingerprintResult("pptx_image", _hash_images(images))
