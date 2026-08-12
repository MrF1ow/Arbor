from __future__ import annotations

import shutil
import subprocess
from pathlib import Path

from pptx import Presentation

from arbor_worker.prepare import PrepareError


def extract_pptx_slide_texts(source: Path) -> list[str]:
    prs = Presentation(str(source))
    slides: list[str] = []
    for slide in prs.slides:
        chunks: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        chunks.append(line)
        slides.append("\n".join(chunks))
    return slides


def extract_pptx_text(source: Path) -> str:
    return "\n".join(t for t in extract_pptx_slide_texts(source) if t.strip())


def find_soffice(which=shutil.which) -> str | None:
    return which("soffice") or which("libreoffice")


def convert_pptx_to_pdf(source: Path, out_dir: Path, runner=subprocess.run, which=shutil.which) -> Path:
    soffice = find_soffice(which)
    if soffice is None:
        raise PrepareError(
            "PPTX has insufficient text and LibreOffice (soffice) is not installed for "
            "image fallback. Re-export the slides as PDF, or install LibreOffice."
        )
    proc = runner(
        [soffice, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(source)],
        capture_output=True,
        text=True,
    )
    if proc.returncode != 0:
        raise PrepareError(f"LibreOffice conversion failed: {(proc.stderr or '').strip()}")
    pdf = out_dir / (source.stem + ".pdf")
    if not pdf.is_file():
        raise PrepareError("LibreOffice did not produce a PDF")
    return pdf
