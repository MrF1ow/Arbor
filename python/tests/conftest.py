from __future__ import annotations

import subprocess
from pathlib import Path

import pytest


def _run_git(root: Path, *args: str) -> str:
    proc = subprocess.run(
        ["git", "-C", str(root), *args],
        capture_output=True,
        text=True,
        check=True,
    )
    return proc.stdout


@pytest.fixture
def git():
    def _factory(root: Path):
        def _call(*args: str) -> str:
            return _run_git(root, *args)
        return _call
    return _factory


@pytest.fixture
def git_repo(tmp_path: Path) -> Path:
    root = tmp_path / "Knowledge"
    root.mkdir()
    _run_git(root, "init", "-q")
    _run_git(root, "config", "user.email", "test@arbor.local")
    _run_git(root, "config", "user.name", "Arbor Test")
    (root / ".gitignore").write_text("_arbor_cache/\n")
    _run_git(root, "add", ".gitignore")
    _run_git(root, "commit", "-q", "-m", "init")
    return root


@pytest.fixture
def make_pdf():
    import fitz  # PyMuPDF

    def _factory(path: Path, pages: int = 2, text: str = "Slide") -> Path:
        doc = fitz.open()
        for i in range(pages):
            page = doc.new_page()
            page.insert_text((72, 72), f"{text} {i + 1}")
        doc.save(str(path))
        doc.close()
        return path

    return _factory


@pytest.fixture
def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches

    def _factory(path: Path, slides_text: list[str]) -> Path:
        prs = Presentation()
        blank = prs.slide_layouts[6]
        for text in slides_text:
            slide = prs.slides.add_slide(blank)
            box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
            box.text_frame.text = text
        prs.save(str(path))
        return path

    return _factory
